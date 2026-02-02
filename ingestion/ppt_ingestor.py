from pathlib import Path
from typing import List
from pptx import Presentation
from langchain_core.documents import Document


class PPTIngestor:
    def __init__(self, ppt_dir: Path):
        self.ppt_dir = ppt_dir

    def ingest(self) -> List[Document]:
        documents: List[Document] = []

        for file in self.ppt_dir.glob("*.pptx"):
            prs = Presentation(file)

            for slide_idx, slide in enumerate(prs.slides, start=1):
                text_blocks: List[str] = []

                # 1️ Extract text frames (titles + bullets)
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        text = shape.text_frame.text.strip()
                        if text:
                            text_blocks.append(text)

                # 2️ Extract tables 
                for shape in slide.shapes:
                    if shape.has_table:
                        for row in shape.table.rows:
                            cells = [
                                cell.text.strip()
                                for cell in row.cells
                                if cell.text.strip()
                            ]
                            if cells:
                                text_blocks.append(" | ".join(cells))

                # 3️ Extract speaker notes 
                if slide.has_notes_slide:
                    notes = slide.notes_slide.notes_text_frame.text.strip()
                    if notes:
                        text_blocks.append(notes)

                # 4️ Build semantic content
                slide_text = "\n".join(text_blocks).strip()

                # 5️ HARD semantic filter 
                if len(slide_text.split()) < 30:
                    continue

                documents.append(
                    Document(
                        page_content=slide_text,
                        metadata={
                            "document_type": "pptx",
                            "source": file.name,
                            "section": f"slide_{slide_idx}",
                        },
                    )
                )

        return documents
